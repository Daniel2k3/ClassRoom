from PyPDF2 import PdfReader
from pptx import Presentation
import os

def extract_text(path):

    # PDF
    if path.endswith(".pdf"):
        reader = PdfReader(path)
        return "\n".join(page.extract_text() for page in reader.pages)

    # PPTX
    if path.endswith(".pptx"):
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return "Unsupported file type."

